#python -m venv ./venv
#pip freeze > requirements.txt
#pip install -r requirements.txt
#streamlit run app.py

import streamlit as st
from dotenv import load_dotenv
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings, HuggingFaceInstructEmbeddings
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
#from files_processing import get_text_from_files
from htmlTemplates import css, bot_template, user_template
from langchain.llms import HuggingFaceHub

from PyPDF2 import PdfReader
import docx
from pptx import Presentation
import streamlit as st
import hashlib
import os

# global variables
# embeddings = OpenAIEmbeddings()
global embeddings, vectorstore, processed_files_set

processed_files_file = "processed_files.txt"

def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks


def get_vectorstore(file_path):
    global embeddings, vectorstore
    if embeddings not in globals() or embeddings is None:
        embeddings = HuggingFaceInstructEmbeddings(model_name="hkunlp/instructor-large")
    if (vectorstore not in globals() or vectorstore is None) and os.path.exists(file_path):
        vectorstore = FAISS.load_local("faiss_index", embeddings)
        print("Loaded vectorstore from file.")
    else:
        vectorstore = None
        print("No vectorstore found.")
    return vectorstore

def add_to_vectorstore(text_chunks):
    #C:\Users\vipul\.cache\torch\sentence_transformers\hkunlp_instructor-large
    global embeddings, vectorstore
    new_vs = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    #vectorstore = new_vs
    if vectorstore is None:
        vectorstore.merge_from(new_vs)
    else:
        vectorstore = new_vs
    vectorstore.save_local("faiss_index")
    return True

def generate_hash(file_path):
    return hashlib.md5(file_path.encode()).hexdigest()

def build_processed_files_list():
    global processed_files_set, processed_files_file
    processed_files_set = set()
    # build dictionary of processed files
    if os.path.exists(processed_files_file):
        with open(processed_files_file, 'r') as f:
            for line in f:
                line = line.strip()
                if line:
                    file_hash, file_path = line.split('|')
                    processed_files_set[file_hash] = file_path
    

def read_dir(directory):
    # process new files
    for root, _, files in os.walk(directory):
        for file in files:
            file_path = os.path.join(root, file)
            file_hash = generate_hash(file_path)
            if file_hash not in processed_files_set:
                add_file_to_vdb(file_path)

def add_file_to_vdb(file_path):
    raw_text = get_text_from_file(file_path)

    # get the text chunks
    text_chunks = get_text_chunks(raw_text)

    # create vector store
    if add_to_vectorstore(text_chunks):
        with open(processed_files_file, 'w') as f:
            file_hash = generate_hash(file_path)
            f.write(f"{file_hash}|{file_path}\n")
        processed_files_set[file_hash] = file_path
    return True

def get_txt_text(txt_file):
    with open(txt_file, 'r', encoding='utf-8') as file:
        text = file.read()
    return text


def get_doc_text(doc_file):
    doc = docx.Document(doc_file)
    text = ""

    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"

    return text


def get_ppt_text(ppt_file):
    prs = Presentation(ppt_file)
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text

    return text


def get_pdf_text(pdf_doc):
    text = ""
    pdf_reader = PdfReader(pdf_doc)
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text


def get_text_from_files(all_files):
    raw_text = ""
    for file in all_files:
        raw_text += get_text_from_file(file)
    return raw_text

def get_text_from_file(file):
    file_extension = file.name.rsplit('.', 1)[-1]
    full_file_with_path = "./files/" + file.name
    if file_extension == 'txt':
        raw_text = get_txt_text(full_file_with_path)
    elif file_extension == 'docx' or file_extension == 'doc':
        raw_text = get_doc_text(full_file_with_path)
    elif file_extension == 'pptx' or file_extension == 'ppt':
        raw_text = get_ppt_text(full_file_with_path)
    elif file_extension == 'pdf':
        raw_text = get_pdf_text(full_file_with_path)
    else:
        st.error("Please upload a file with one of the following extensions: txt, docx, pptx, pdf")
    return raw_text

def get_conversation_chain(vectorstore):
    llm = ChatOpenAI()
    # llm = HuggingFaceHub(repo_id="google/flan-t5-xxl", model_kwargs={"temperature":0.5, "max_length":512})

    memory = ConversationBufferMemory(
        memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain

def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']

    for i, message in enumerate(st.session_state.chat_history):
        if i % 2 == 0:
            st.write(user_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)
        else:
            st.write(bot_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)

def main():
    global vectorstore, embeddings
    vectorstore = None
    embeddings = None
    build_processed_files_list()
    get_vectorstore("faiss_index\index.faiss")
    load_dotenv()
    st.set_page_config(page_title="Chat with multiple Documents",
                       page_icon=":books:")
    st.write(css, unsafe_allow_html=True)

    if "conversation" not in st.session_state:
        st.session_state.conversation = None
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None

    # create conversation chain
    st.session_state.conversation = get_conversation_chain(vectorstore)

    st.header("Chat with multiple Documents :books:")
    user_question = st.text_input("Ask a question about your documents:")
    if user_question:
        handle_userinput(user_question)

    with st.sidebar:
        st.subheader("Your documents")
        directory = st.file_uploader("Select a directory", type="directory")
        all_files = st.file_uploader(
            "Upload your Documents here and click on 'Process'", accept_multiple_files=True)
        if st.button("Process"):
            with st.spinner("Processing"):
                if all_files:
                    raw_text = get_text_from_files(all_files)
                    text_chunks = get_text_chunks(raw_text)
                    vectorstore = add_to_vectorstore(text_chunks)
                if directory:
                    directory_to_parse = directory.name
                    processed_files_file = 'processed_files.txt'
                    read_dir(directory_to_parse, processed_files_file)


if __name__ == '__main__':
    main()
